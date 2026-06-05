#!/usr/bin/env python3
"""
Two-pass document categorization for /Volumes/toshiba.

Pass 1:  python categorize.py sample
Pass 2:  python categorize.py run [--dry-run]
"""
import argparse
import json
import math
import os
import random
import re
import shutil
import subprocess
import sys
from pathlib import Path

import openai
from dotenv import load_dotenv

load_dotenv()

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

DRIVE_PATH = Path("/Volumes/toshiba")
ORGANIZED_PATH = DRIVE_PATH / "_Organized"

DOCUMENT_EXTENSIONS = {
    ".pdf", ".doc", ".docx", ".txt", ".md", ".rtf",
    ".csv", ".xlsx", ".xls", ".pptx", ".ppt",
    ".pages", ".numbers", ".keynote", ".odt", ".ods", ".odp",
}

SAMPLE_SIZE = 200
SNIPPET_CHARS = 400   # chars extracted per file for sampling / classification
BATCH_SIZE = 20       # files per API batch

SCRIPT_DIR = Path(__file__).parent
CATEGORIES_FILE = SCRIPT_DIR / "categories.txt"
PROGRESS_FILE = SCRIPT_DIR / "progress.json"

MODEL = "gpt-4o-mini"


# ---------------------------------------------------------------------------
# Guardian
# ---------------------------------------------------------------------------

class RunGuardian:
    """Monitors token spend, unsorted rate, and category skew. Pauses on breach."""

    def __init__(
        self,
        token_budget: int,
        max_unsorted_rate: float,
        max_skew_rate: float,
        initial_tokens: int = 0,
    ) -> None:
        self.token_budget = token_budget
        self.max_unsorted_rate = max_unsorted_rate
        self.max_skew_rate = max_skew_rate
        self.tokens_used = initial_tokens
        self._category_counts: dict[str, int] = {}

    @property
    def total_classified(self) -> int:
        return sum(self._category_counts.values())

    def record_usage(self, batch_tokens: int) -> None:
        self.tokens_used += batch_tokens

    def record_batch(self, classifications: dict[str, str]) -> None:
        for cat in classifications.values():
            self._category_counts[cat] = self._category_counts.get(cat, 0) + 1

    def check(self) -> tuple[bool, str]:
        """Return (ok, reason). ok=False means pause the run."""
        if self.tokens_used >= self.token_budget:
            return False, (
                f"token budget reached: {self.tokens_used:,}/{self.token_budget:,} tokens used"
            )

        if self.total_classified >= 60:  # wait for 3 batches before judging
            unsorted = self._category_counts.get("_Unsorted", 0)
            unsorted_rate = unsorted / self.total_classified
            if unsorted_rate > self.max_unsorted_rate:
                return False, (
                    f"_Unsorted rate {unsorted_rate:.1%} exceeds limit {self.max_unsorted_rate:.1%} "
                    f"— categories may not match file content"
                )

            non_unsorted = {k: v for k, v in self._category_counts.items() if k != "_Unsorted"}
            if non_unsorted and unsorted > 0:
                max_cat = max(non_unsorted, key=non_unsorted.__getitem__)
                skew = non_unsorted[max_cat] / self.total_classified
                if skew > self.max_skew_rate:
                    return False, (
                        f"category '{max_cat}' has {skew:.1%} of files "
                        f"(limit {self.max_skew_rate:.1%}) — taxonomy may be too narrow"
                    )

        return True, ""

    def print_stats(self) -> None:
        est_cost = (self.tokens_used * 0.225) / 1_000_000  # blended ~75% input / 25% output
        print(f"\n  Tokens used : {self.tokens_used:,} / {self.token_budget:,}")
        print(f"  Est. cost   : ${est_cost:.4f}")
        if self._category_counts:
            print(f"  Distribution ({self.total_classified} files):")
            for cat, count in sorted(self._category_counts.items(), key=lambda x: -x[1]):
                pct = count / self.total_classified * 100
                print(f"    {cat:<22} {count:>5}  ({pct:.1f}%)")


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext in (".txt", ".md", ".csv"):
            return path.read_text(errors="ignore")[:SNIPPET_CHARS]

        if ext == ".pdf":
            try:
                import pypdf
                reader = pypdf.PdfReader(str(path))
                text = ""
                for page in reader.pages[:3]:
                    text += page.extract_text() or ""
                    if len(text) >= SNIPPET_CHARS:
                        break
                return text[:SNIPPET_CHARS]
            except Exception:
                return ""

        if ext in (".doc", ".docx"):
            try:
                from docx import Document
                doc = Document(str(path))
                text = " ".join(p.text for p in doc.paragraphs[:30])
                return text[:SNIPPET_CHARS]
            except Exception:
                return ""

        if ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
                ws = wb.active
                cells = []
                for row in ws.iter_rows(max_row=15, values_only=True):
                    cells.extend(str(c) for c in row if c is not None)
                return " ".join(cells)[:SNIPPET_CHARS]
            except Exception:
                return ""

        if ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                text = ""
                for slide in prs.slides[:5]:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + " "
                return text[:SNIPPET_CHARS]
            except Exception:
                return ""

        if ext == ".rtf":
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", str(path)],
                capture_output=True, text=True, timeout=10,
            )
            return result.stdout[:SNIPPET_CHARS]

    except Exception:
        pass

    return ""


def is_document(path: Path) -> bool:
    return path.suffix.lower() in DOCUMENT_EXTENSIONS


# ---------------------------------------------------------------------------
# File walking
# ---------------------------------------------------------------------------

def walk_files(root: Path, skip_dirs: set[str] | None = None) -> list[Path]:
    skip_dirs = skip_dirs or set()
    files = []
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [
            d for d in dirnames
            if d not in skip_dirs and not d.startswith(".")
        ]
        for name in filenames:
            if not name.startswith("."):
                files.append(Path(dirpath) / name)
    return files


# ---------------------------------------------------------------------------
# Destination helpers
# ---------------------------------------------------------------------------

def unique_dest(dest: Path) -> Path:
    if not dest.exists():
        return dest
    counter = 1
    while True:
        candidate = dest.parent / f"{dest.stem}_{counter}{dest.suffix}"
        if not candidate.exists():
            return candidate
        counter += 1


# ---------------------------------------------------------------------------
# Pre-flight estimate
# ---------------------------------------------------------------------------

def pre_flight_estimate(
    doc_files: list[Path],
    other_files: list[Path],
    guardian: RunGuardian,
) -> bool:
    """Print cost/token estimate and ask user to confirm. Returns True to proceed."""
    batch_count = math.ceil(len(doc_files) / BATCH_SIZE) if doc_files else 0
    est_input_tokens = int(len(doc_files) * SNIPPET_CHARS / 4) + batch_count * 200
    est_output_tokens = batch_count * 200
    est_total = est_input_tokens + est_output_tokens
    # gpt-4o-mini: $0.15/1M input, $0.60/1M output
    est_cost = (est_input_tokens * 0.15 + est_output_tokens * 0.60) / 1_000_000
    tokens_remaining = guardian.token_budget - guardian.tokens_used

    print("Pre-flight estimate")
    print(f"  Documents to classify : {len(doc_files):,}")
    print(f"  Other files (_Other)  : {len(other_files):,}")
    print(f"  API batches           : {batch_count:,}")
    print(f"  Est. tokens           : ~{est_total:,}")
    print(f"  Est. cost             : ~${est_cost:.4f}")
    print(f"  Budget remaining      : {tokens_remaining:,} tokens")
    if est_total > tokens_remaining:
        print()
        print("  Warning: estimated usage exceeds remaining budget.")
        print("  The run will process files until the budget is reached, then pause.")
        print("  Raise TOKEN_BUDGET in .env to increase the cap.")
    print()
    answer = input("Proceed? [y/N] ").strip().lower()
    return answer == "y"


# ---------------------------------------------------------------------------
# Progress tracking (so the run is resumable)
# ---------------------------------------------------------------------------

def load_progress() -> tuple[set[str], int]:
    """Returns (processed_paths, cumulative_tokens_used)."""
    if PROGRESS_FILE.exists():
        data = json.loads(PROGRESS_FILE.read_text())
        if isinstance(data, list):
            # migrate old format (plain list of paths)
            return set(data), 0
        return set(data.get("processed", [])), int(data.get("tokens_used", 0))
    return set(), 0


def save_progress(processed: set[str], tokens_used: int) -> None:
    PROGRESS_FILE.write_text(json.dumps({
        "processed": list(processed),
        "tokens_used": tokens_used,
    }))


# ---------------------------------------------------------------------------
# OpenAI helpers
# ---------------------------------------------------------------------------

def parse_json_array(text: str) -> list[str]:
    text = text.strip()
    try:
        result = json.loads(text)
        if isinstance(result, list):
            return [str(x) for x in result]
    except json.JSONDecodeError:
        pass
    return re.findall(r'"([^"]+)"', text)


def parse_json_object(text: str) -> dict[str, str]:
    text = text.strip()
    try:
        result = json.loads(text)
        if isinstance(result, dict):
            return {str(k): str(v) for k, v in result.items()}
    except json.JSONDecodeError:
        pass
    return {}


# ---------------------------------------------------------------------------
# Pass 1 — sample
# ---------------------------------------------------------------------------

def sample_command() -> None:
    client = openai.OpenAI()

    print(f"Walking {DRIVE_PATH} ...")
    all_files = walk_files(DRIVE_PATH, skip_dirs={"_Organized", "_deleted"})
    doc_files = [f for f in all_files if is_document(f)]
    print(f"Found {len(doc_files)} document files. Sampling {min(SAMPLE_SIZE, len(doc_files))}.")

    sample = random.sample(doc_files, min(SAMPLE_SIZE, len(doc_files)))

    snippets = []
    for path in sample:
        text = extract_text(path)
        if text.strip():
            snippets.append(f"File: {path.name}\nContent: {text[:300]}")

    print(f"Extracted text from {len(snippets)} files. Sending to model in batches...")

    raw_categories: list[str] = []
    total_batches = (len(snippets) + BATCH_SIZE - 1) // BATCH_SIZE

    for i in range(0, len(snippets), BATCH_SIZE):
        batch = snippets[i : i + BATCH_SIZE]
        batch_num = i // BATCH_SIZE + 1
        print(f"  Batch {batch_num}/{total_batches} ...")

        response = client.chat.completions.create(
            model=MODEL,
            max_tokens=400,
            messages=[{
                "role": "user",
                "content": (
                    "You are analyzing a sample of files from a personal archive to suggest a category taxonomy.\n\n"
                    f"Here are {len(batch)} file samples:\n\n"
                    + "\n\n---\n\n".join(batch)
                    + "\n\nBased on these files, suggest broad category names suitable for organizing this archive. "
                    "Reply with ONLY a JSON array of category name strings, nothing else. "
                    'Example: ["Finance", "Contracts", "Personal", "Work"]'
                ),
            }],
        )
        raw_categories.extend(parse_json_array(response.choices[0].message.content))

    print("\nConsolidating categories...")
    response = client.chat.completions.create(
        model=MODEL,
        max_tokens=300,
        messages=[{
            "role": "user",
            "content": (
                "Here is a raw list of document category names from multiple batches. "
                "Some may be duplicates or near-duplicates.\n\n"
                + json.dumps(list(set(raw_categories)))
                + "\n\nConsolidate into a clean final list of 5-15 broad categories "
                "suitable for a personal file archive. Reply with ONLY a JSON array of strings."
            ),
        }],
    )
    final_categories = parse_json_array(response.choices[0].message.content)

    if not final_categories:
        final_categories = ["Finance", "Personal", "Work", "Legal", "Medical", "Travel"]

    lines = [
        "# Edit this list before running: python categorize.py run",
        "# One category per line. Lines starting with # are ignored.",
        "",
    ] + final_categories

    CATEGORIES_FILE.write_text("\n".join(lines) + "\n")

    print(f"\nProposed categories written to: {CATEGORIES_FILE}")
    print("\nCategories:")
    for cat in final_categories:
        print(f"  - {cat}")
    print(f"\nEdit {CATEGORIES_FILE.name} if needed, then run:")
    print("  python categorize.py run --dry-run")


# ---------------------------------------------------------------------------
# Pass 2 — run
# ---------------------------------------------------------------------------

def load_categories() -> list[str]:
    if not CATEGORIES_FILE.exists():
        sys.exit(f"Error: {CATEGORIES_FILE} not found. Run 'python categorize.py sample' first.")
    categories = [
        line.strip()
        for line in CATEGORIES_FILE.read_text().splitlines()
        if line.strip() and not line.startswith("#")
    ]
    if not categories:
        sys.exit(f"Error: No categories found in {CATEGORIES_FILE}.")
    return categories


def classify_batch(
    client: openai.OpenAI,
    batch: list[tuple[Path, str]],
    categories: list[str],
) -> tuple[dict[str, str], int]:
    """Returns (classifications, total_tokens_used)."""
    cat_list = ", ".join(f'"{c}"' for c in categories)
    file_blocks = [
        f"File: {path.name}\nContent: {(snippet or '(no text)')[:300]}"
        for path, snippet in batch
    ]
    prompt = (
        f"Classify each file into exactly one of these categories: {cat_list}\n\n"
        'If a file does not fit any category, use "_Unsorted".\n\n'
        "Files to classify:\n\n"
        + "\n\n---\n\n".join(file_blocks)
        + "\n\nReply with ONLY a JSON object mapping each filename to its category. "
        'Example: {"report.pdf": "Finance", "letter.docx": "Personal"}'
    )

    response = client.chat.completions.create(
        model=MODEL,
        max_tokens=800,
        messages=[{"role": "user", "content": prompt}],
    )
    result = parse_json_object(response.choices[0].message.content)
    if not result:
        result = {path.name: "_Unsorted" for path, _ in batch}
    return result, response.usage.total_tokens


def run_command(dry_run: bool) -> None:
    categories = load_categories()
    valid_categories = set(categories) | {"_Unsorted", "_Other"}
    client = openai.OpenAI()
    processed, tokens_used = load_progress()

    token_budget = int(os.getenv("TOKEN_BUDGET", "2000000"))
    max_unsorted_rate = float(os.getenv("MAX_UNSORTED_RATE", "0.50"))
    max_skew_rate = float(os.getenv("MAX_SKEW_RATE", "0.80"))
    guardian = RunGuardian(token_budget, max_unsorted_rate, max_skew_rate, initial_tokens=tokens_used)

    label = "DRY RUN — no files will be moved" if dry_run else "LIVE RUN — files will be moved"
    print(f"{label}")
    print(f"Categories : {', '.join(categories)}")
    print(f"Output     : {ORGANIZED_PATH}\n")

    all_files = walk_files(DRIVE_PATH, skip_dirs={"_Organized", "_deleted"})
    remaining = [f for f in all_files if str(f) not in processed]
    print(f"Total files: {len(all_files)}  |  Already processed: {len(processed)}  |  Remaining: {len(remaining)}\n")

    doc_files = [f for f in remaining if is_document(f)]
    other_files = [f for f in remaining if not is_document(f)]

    if not dry_run:
        if not pre_flight_estimate(doc_files, other_files, guardian):
            print("Aborted.")
            return
        for cat in categories:
            (ORGANIZED_PATH / cat).mkdir(parents=True, exist_ok=True)
        (ORGANIZED_PATH / "_Unsorted").mkdir(parents=True, exist_ok=True)
        (ORGANIZED_PATH / "_Other").mkdir(parents=True, exist_ok=True)

    moved = errors = 0

    # Non-documents → _Other
    for path in other_files:
        dest = ORGANIZED_PATH / "_Other" / path.name
        if dry_run:
            print(f"  [_Other]       {path.relative_to(DRIVE_PATH)}")
        else:
            try:
                shutil.move(str(path), str(unique_dest(dest)))
                processed.add(str(path))
                moved += 1
            except Exception as e:
                print(f"  ERROR {path}: {e}")
                errors += 1

    # Documents → classified
    total_batches = (len(doc_files) + BATCH_SIZE - 1) // BATCH_SIZE
    batch: list[tuple[Path, str]] = []

    for idx, path in enumerate(doc_files):
        batch.append((path, extract_text(path)))

        if len(batch) == BATCH_SIZE or idx == len(doc_files) - 1:
            batch_num = idx // BATCH_SIZE + 1
            print(f"Classifying batch {batch_num}/{total_batches} ({len(batch)} files)...")

            if dry_run:
                classifications = {path.name: "(dry-run)" for path, _ in batch}
                batch_tokens = 0
            else:
                try:
                    classifications, batch_tokens = classify_batch(client, batch, categories)
                except Exception as e:
                    print(f"  API error: {e} — marking batch as _Unsorted")
                    classifications = {p.name: "_Unsorted" for p, _ in batch}
                    batch_tokens = 0

            if not dry_run:
                guardian.record_usage(batch_tokens)
                guardian.record_batch(classifications)

            for path, _ in batch:
                category = classifications.get(path.name, "_Unsorted")
                if category not in valid_categories:
                    category = "_Unsorted"

                dest = ORGANIZED_PATH / category / path.name
                if dry_run:
                    print(f"  [{category:<20}] {path.relative_to(DRIVE_PATH)}")
                else:
                    try:
                        shutil.move(str(path), str(unique_dest(dest)))
                        processed.add(str(path))
                        moved += 1
                    except Exception as e:
                        print(f"  ERROR {path}: {e}")
                        errors += 1

            if not dry_run:
                save_progress(processed, guardian.tokens_used)

                ok, reason = guardian.check()
                if not ok:
                    print(f"\nGuardian: pausing run — {reason}")
                    guardian.print_stats()
                    print(f"\nProgress saved to {PROGRESS_FILE.name}.")
                    print("Resume with:\n  python categorize.py run")
                    sys.exit(0)

            batch = []

    summary = "Would move" if dry_run else "Moved"
    count = len(other_files) + len(doc_files) if dry_run else moved
    print(f"\n{summary}: {count} files | Errors: {errors}")

    if not dry_run and errors == 0:
        PROGRESS_FILE.unlink(missing_ok=True)
        print("Progress file cleaned up.")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Two-pass document categorization for /Volumes/toshiba"
    )
    sub = parser.add_subparsers(dest="command", required=True)

    sub.add_parser("sample", help="Pass 1: sample files and discover categories")

    run_p = sub.add_parser("run", help="Pass 2: classify and move files")
    run_p.add_argument(
        "--dry-run", action="store_true",
        help="Preview moves without touching any files",
    )

    args = parser.parse_args()

    if args.command == "sample":
        sample_command()
    elif args.command == "run":
        run_command(dry_run=args.dry_run)


if __name__ == "__main__":
    main()
